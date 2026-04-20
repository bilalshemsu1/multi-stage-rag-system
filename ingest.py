from __future__ import annotations

import argparse
import hashlib
import os
import re
import sys
import time
import socket
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import fitz  # PyMuPDF
from rich.console import Console
from rich.panel import Panel
from rich.progress import (
    Progress,
    SpinnerColumn,
    TextColumn,
    BarColumn,
    TaskProgressColumn,
    TimeElapsedColumn,
    TimeRemainingColumn,
)
from rich.table import Table
from rich.live import Live
from rich.layout import Layout
from rich import box

from core import (
    DEFAULT_SOURCE_DIR,
    DEFAULT_CHROMA_DIR,
    DEFAULT_COLLECTION,
    DEFAULT_EMBED_BATCH,
    DEFAULT_EMBED_MODEL,
    ChromaVectorStore,
    SentenceTransformerEmbedder,
)

# Chunking tuned for recall (more text per result) while staying reasonably sized for embeddings
DEFAULT_CHUNK_CHARS = int(os.environ.get("CHUNK_CHARS", "2200"))
DEFAULT_CHUNK_OVERLAP_CHARS = int(os.environ.get("CHUNK_OVERLAP_CHARS", "350"))
DEFAULT_MIN_CHUNK_CHARS = int(os.environ.get("MIN_CHUNK_CHARS", "120"))
DEFAULT_BATCH_PAGES = 20 # Number of pages to process at once for large files

console = Console()

def _now() -> float:
    return time.time()

def _sha1_text(s: str) -> str:
    return hashlib.sha1(s.encode("utf-8", errors="ignore")).hexdigest()

def _sha1_file(path: Path, max_bytes: int = 8 * 1024 * 1024) -> str:
    h = hashlib.sha1()
    with path.open("rb") as f:
        remaining = max_bytes
        while remaining > 0:
            chunk = f.read(min(1024 * 1024, remaining))
            if not chunk:
                break
            h.update(chunk)
            remaining -= len(chunk)
    return h.hexdigest()

def _normalize_ws(text: str) -> str:
    return re.sub(r"\s+", " ", text).strip()

def _safe_int(x: Any, default: int = 0) -> int:
    try:
        return int(x)
    except Exception:
        return default

def _extract_subject(document_name: str) -> str:
    name_lower = document_name.lower()
    subject_mapping = {
        "english": ["english", "eng"],
        "mathematics": ["mathematics", "math", "maths"],
        "biology": ["biology", "bio"],
        "physics": ["physics"],
        "chemistry": ["chemistry", "chem"],
        "history": ["history", "hist"],
        "geography": ["geography", "geo"],
        "civics": ["civics", "citizenship"],
        "economics": ["economics", "econ"],
        "business": ["business", "commerce"],
        "ethics": ["ethics", "moral"],
        "science": ["science", "sci"],
        "language": ["language", "grammar"],
        "social": ["social studies", "social"],
        "computer": ["computer", "ict", "it"],
        "health": ["health", "pe", "physical"],
    }
    for subject, keywords in subject_mapping.items():
        if any(kw in name_lower for kw in keywords):
            return subject
    words = re.split(r"[\s_\-]+", Path(document_name).stem)
    return (words[0].lower() if words and words[0] else "other")[:64]

def is_internet_available(host="8.8.8.8", port=53, timeout=3):
    """Check if the internet is accessible."""
    try:
        socket.setdefaulttimeout(timeout)
        socket.socket(socket.AF_INET, socket.SOCK_STREAM).connect((host, port))
        return True
    except socket.error:
        return False

@dataclass(frozen=True)
class Chunk:
    text: str
    start_page: int
    end_page: int
    heading: str
    chunk_index: int
    total_chunks: int

class DocumentChunker:
    def __init__(self, chunk_chars: int, overlap_chars: int, min_chunk_chars: int):
        self.chunk_chars = chunk_chars
        self.overlap_chars = overlap_chars
        self.min_chunk_chars = min_chunk_chars

    def extract_elements(self, file_path: Path, include_tables: bool, include_images_markers: bool, start_page: int = 0, end_page: Optional[int] = None) -> Tuple[List[Dict[str, Any]], Dict[str, int]]:
        ext = file_path.suffix.lower()
        if ext == ".pdf":
            return self._extract_pdf(file_path, include_tables, include_images_markers, start_page, end_page)
        elif ext == ".docx":
            return self._extract_docx(file_path)
        elif ext == ".pptx":
            return self._extract_pptx(file_path)
        else:
            return [], {"pages": 0, "tables": 0, "images": 0}

    def _extract_docx(self, path: Path) -> Tuple[List[Dict[str, Any]], Dict[str, int]]:
        try:
            import docx
            doc = docx.Document(path)
        except Exception as e:
            return [], {"pages": 0, "tables": 0, "images": 0}

        elements: List[Dict[str, Any]] = []
        for para in doc.paragraphs:
            txt = para.text.strip()
            if not txt:
                continue
            
            is_heading = False
            style_name = str(para.style.name if para.style else "").lower()
            if "heading" in style_name or "title" in style_name:
                is_heading = True
            elif para.runs and para.runs[0].bold and len(txt) < 80:
                is_heading = True

            elements.append({
                "type": "heading" if is_heading else "text",
                "content": txt,
                "page": 1,
                "bbox": [0, 0, 0, 0],
            })
        
        return elements, {"pages": 1, "tables": 0, "images": 0}

    def _extract_pptx(self, path: Path) -> Tuple[List[Dict[str, Any]], Dict[str, int]]:
        try:
            from pptx import Presentation
            prs = Presentation(path)
        except Exception as e:
            return [], {"pages": 0, "tables": 0, "images": 0}

        elements: List[Dict[str, Any]] = []
        for slide_idx, slide in enumerate(prs.slides):
            page_no = slide_idx + 1
            for shape_idx, shape in enumerate(slide.shapes):
                if not shape.has_text_frame:
                    continue
                txt = shape.text.strip()
                if not txt:
                    continue
                is_heading = (shape_idx == 0)
                elements.append({
                    "type": "heading" if is_heading else "text",
                    "content": txt,
                    "page": page_no,
                    "bbox": [0, 0, 0, 0],
                })
        return elements, {"pages": len(prs.slides), "tables": 0, "images": 0}

    def _extract_pdf(self, pdf_path: Path, include_tables: bool, include_images_markers: bool, start_page: int = 0, end_page: Optional[int] = None) -> Tuple[List[Dict[str, Any]], Dict[str, int]]:
        doc = fitz.open(pdf_path)
        total_len = len(doc)
        elements: List[Dict[str, Any]] = []
        counters = {"pages": 0, "tables": 0, "images": 0}
        
        target_end = end_page if end_page is not None else total_len
        target_end = min(target_end, total_len)
        
        try:
            for page_idx in range(start_page, target_end):
                page = doc[page_idx]
                page_no = page_idx + 1
                page_elements: List[Dict[str, Any]] = []
                try:
                    blocks = page.get_text("dict").get("blocks", [])
                    for block in blocks:
                        if "lines" not in block: continue
                        for line in block["lines"]:
                            for span in line.get("spans", []):
                                txt = (span.get("text") or "").strip()
                                if len(txt) < 2: continue
                                font_size = float(span.get("size") or 0)
                                font_name = str(span.get("font") or "").lower()
                                is_heading = bool(
                                    font_size >= 14.5 or "bold" in font_name
                                    or (len(txt) <= 80 and txt.isupper())
                                    or re.match(r"^(chapter|section|unit|part)\s+\d+", txt, re.IGNORECASE)
                                )
                                page_elements.append({
                                    "type": "heading" if is_heading else "text",
                                    "content": txt, "page": page_no, "bbox": span.get("bbox", [0, 0, 0, 0])
                                })
                except Exception: pass
                if include_tables:
                    try:
                        tables = page.find_tables()
                        for table in tables:
                            md = self._table_to_markdown(table.extract())
                            if md and md != "[Empty table]":
                                page_elements.append({
                                    "type": "table", "content": f"[TABLE]\n{md}", "page": page_no, "bbox": getattr(table, "bbox", [0, 0, 0, 0])
                                })
                                counters["tables"] += 1
                    except Exception: pass
                if include_images_markers:
                    try:
                        imgs = page.get_images(full=True) or []
                        for img_idx, _ in enumerate(imgs):
                            page_elements.append({
                                "type": "image", "content": f"[IMAGE {img_idx+1}]", "page": page_no, "bbox": [0, 0, 0, 0]
                            })
                            counters["images"] += 1
                    except Exception: pass
                page_elements.sort(key=lambda x: (x["bbox"][1], x["bbox"][0]))
                elements.extend(page_elements)
                counters["pages"] += 1
        finally:
            doc.close()
        return elements, counters

    @staticmethod
    def _table_to_markdown(table_data: Any) -> str:
        if not table_data: return "[Empty table]"
        cleaned = []
        for row in table_data:
            if row is None: continue
            cleaned_row = [str(cell).strip().replace("\n", " ") if cell is not None else "" for cell in row]
            if any(c for c in cleaned_row): cleaned.append(cleaned_row)
        if not cleaned: return "[Empty table]"
        max_cols = max(len(r) for r in cleaned)
        for r in cleaned:
            while len(r) < max_cols: r.append("")
        lines = [" | ".join(cleaned[0]), " | ".join(["---"] * max_cols)]
        for r in cleaned[1:]: lines.append(" | ".join(r))
        return "\n".join(lines)

    def chunk(self, elements: List[Dict[str, Any]]) -> List[Tuple[str, int, int, str]]:
        if not elements: return []
        chunks: List[Tuple[str, int, int, str]] = []
        buf: List[str] = []
        heading = ""
        start_page = _safe_int(elements[0].get("page"), 1)
        end_page = start_page
        def flush() -> None:
            nonlocal buf, start_page, end_page
            if not buf: return
            body = _normalize_ws(" ".join(buf))
            text = f"## {heading}\n\n{body}".strip() if heading else body
            if len(text) >= self.min_chunk_chars:
                chunks.append((text, start_page, end_page, heading[:200]))
            tail = " ".join(buf)[-self.overlap_chars :] if self.overlap_chars > 0 else ""
            buf = [tail] if tail else []
            start_page = end_page
        cur_len = 0
        for el in elements:
            t = (el.get("content") or "").strip()
            if not t: continue
            page = _safe_int(el.get("page"), end_page)
            if el.get("type") == "heading":
                if cur_len >= int(self.chunk_chars * 0.65):
                    flush()
                    cur_len = len(" ".join(buf))
                heading = t
            if page > end_page and cur_len >= int(self.chunk_chars * 0.55):
                flush()
                cur_len = len(" ".join(buf))
                start_page = page
            if cur_len + len(t) + 1 > self.chunk_chars and cur_len >= int(self.chunk_chars * 0.4):
                flush()
                cur_len = len(" ".join(buf))
                start_page = page
            buf.append(t)
            cur_len += len(t) + 1
            end_page = page
        flush()
        return chunks

class RagDocumentProcessor:
    def __init__(self, source_dir: Path, chroma_dir: Path, collection: str, force_reprocess: bool, include_tables: bool, include_images_markers: bool, chunk_chars: int, overlap_chars: int, min_chunk_chars: int, embed_model: str, embed_batch: int):
        self.source_dir = Path(source_dir)
        self.force_reprocess = force_reprocess
        self.include_tables = include_tables
        self.include_images_markers = include_images_markers
        self.embed_model_name = embed_model
        self.embed_batch = embed_batch
        self.chroma_dir = chroma_dir
        self.collection_name = collection
        self.chunker = DocumentChunker(chunk_chars, overlap_chars, min_chunk_chars)
        self.stats = {"start_time": _now(), "total_files": 0, "processed": 0, "skipped": 0, "failed": 0, "pages": 0, "tables": 0, "images": 0, "chunks_added": 0, "embedding_time_s": 0.0, "store_time_s": 0.0}
        self.embedder = None
        self.store = None

    def initialize(self):
        with console.status("[bold blue]Checking connection & models..."):
            is_connected = is_internet_available()
            if not is_connected:
                console.print("[yellow]⚠ Internet disconnected. Operating in offline mode.")
            
            try:
                self.embedder = SentenceTransformerEmbedder(self.embed_model_name, self.embed_batch)
                self.store = ChromaVectorStore(self.chroma_dir, self.collection_name, embedding_dim_hint=self.embedder.dim, embedder=self.embedder)
            except Exception as e:
                console.print(f"[bold red]Initialization Failed:[/bold red] {e}")
                if not is_connected:
                    console.print("[red]The model may need an initial internet connection to download.")
                sys.exit(1)

    def _doc_processed_marker_id(self, document_name: str, file_sha1: str) -> str:
        return f"doc::{_sha1_text(document_name)}::{file_sha1[:16]}"

    def _is_processed(self, document_name: str, file_sha1: str) -> bool:
        if self.force_reprocess: return False
        marker_id = self._doc_processed_marker_id(document_name, file_sha1)
        try:
            res = self.store.collection.get(ids=[marker_id])
            return bool(res.get("ids"))
        except Exception: return False

    def process_file(self, file_path: Path, progress: Progress, task_id: Any) -> bool:
        doc_name = file_path.name
        try:
            progress.update(task_id, description=f"[cyan]Hashing {doc_name}")
            file_sha1 = _sha1_file(file_path)
            if self._is_processed(doc_name, file_sha1):
                self.stats["skipped"] += 1
                return True
            
            if self.force_reprocess:
                ids = self.store.get_ids_by_doc(doc_name)
                self.store.delete_ids(ids)

            # Pre-flight to get total pages
            ext = file_path.suffix.lower()
            total_pages = 0
            if ext == ".pdf":
                with fitz.open(file_path) as d: total_pages = len(d)
            else:
                total_pages = 1 # Fallback for now

            subject = _extract_subject(doc_name)
            processed_any = False
            
            # Process in page-batches
            for start in range(0, total_pages, DEFAULT_BATCH_PAGES):
                end = min(start + DEFAULT_BATCH_PAGES, total_pages)
                batch_label = f"(Pages {start+1}-{end})" if total_pages > DEFAULT_BATCH_PAGES else ""
                
                progress.update(task_id, description=f"[cyan]Extracting {doc_name} {batch_label}")
                elements, counters = self.chunker.extract_elements(file_path, self.include_tables, self.include_images_markers, start, end)
                self.stats["pages"] += counters["pages"]
                self.stats["tables"] += counters["tables"]
                self.stats["images"] += counters["images"]

                if not elements: continue

                raw_chunks = self.chunker.chunk(elements)
                if not raw_chunks: continue

                chunk_ids, docs, metas = [], [], []
                for idx, (text, start_p, end_p, heading) in enumerate(raw_chunks):
                    norm = _normalize_ws(text)
                    if len(norm) < self.chunker.min_chunk_chars: continue
                    chunk_hash = _sha1_text(norm)[:20]
                    # Unique ID includes the page range to avoid collisions between batches
                    chunk_id = f"chunk::{_sha1_text(doc_name)[:10]}::{start}_{chunk_hash}"
                    metas.append({
                        "type": "chunk", "document_name": doc_name, "subject": subject,
                        "source_path": str(file_path), "file_sha1": file_sha1, "chunk_hash": chunk_hash,
                        "start_page": int(start_p), "end_page": int(end_p), "heading": heading,
                        "processed_at": _now()
                    })
                    chunk_ids.append(chunk_id); docs.append(norm)

                if docs:
                    progress.update(task_id, description=f"[cyan]Embedding {doc_name} {batch_label}")
                    t0 = _now()
                    vectors = self.embedder.embed(docs)
                    self.stats["embedding_time_s"] += _now() - t0

                    progress.update(task_id, description=f"[cyan]Storing {doc_name} {batch_label}")
                    t1 = _now()
                    self.store.upsert(ids=chunk_ids, documents=docs, metadatas=metas, embeddings=vectors)
                    self.stats["store_time_s"] += _now() - t1
                    self.stats["chunks_added"] += len(docs)
                    processed_any = True

            if processed_any:
                # Add doc marker once fully finished
                self.store.upsert(
                    ids=[self._doc_processed_marker_id(doc_name, file_sha1)],
                    documents=[f"[DOC_MARKER] {doc_name}"],
                    metadatas=[{"type": "doc_marker", "document_name": doc_name, "file_sha1": file_sha1}],
                    embeddings=[self.embedder.embed(["marker"])[0]]
                )
                self.stats["processed"] += 1
                return True
            else:
                self.stats["failed"] += 1
                return False

        except Exception as e:
            self.stats["failed"] += 1
            return False

    def make_stats_table(self):
        table = Table(box=box.MINIMAL_DOUBLE_HEAD, expand=True)
        table.add_column("Metric", style="bold cyan")
        table.add_column("Value", justify="right")
        table.add_row("Files Processed", f"[bold green]{self.stats['processed']}")
        table.add_row("Files Skipped", f"[yellow]{self.stats['skipped']}")
        table.add_row("Files Failed", f"[red]{self.stats['failed']}")
        table.add_row("Total Pages", str(self.stats["pages"]))
        table.add_row("Chunks Created", f"[bold blue]{self.stats['chunks_added']:,}")
        
        # Hardware Indicator
        hw_info = "[bold blue]CPU"
        if self.embedder and hasattr(self.embedder.model, "device"):
            dev = str(self.embedder.model.device)
            if "xpu" in dev: hw_info = "[bold green]INTEL GPU (XPU)"
            elif "cuda" in dev: hw_info = "[bold green]NVIDIA GPU (CUDA)"
            elif "mps" in dev: hw_info = "[bold green]APPLE M1 (MPS)"
        table.add_row("Hardware Accelerator", hw_info)
        
        return table

    def process_all(self):
        files = []
        for ext in ["*.pdf", "*.docx", "*.pptx"]:
            files.extend([p for p in self.source_dir.rglob(ext) if not p.name.startswith(".")])
        if not files:
            console.print("[bold red]No documents found in source directory.[/bold red]")
            return

        self.stats["total_files"] = len(files)
        
        # UI Setup
        progress = Progress(
            SpinnerColumn(),
            TextColumn("[progress.description]{task.description}"),
            BarColumn(bar_width=None),
            TaskProgressColumn(),
            TimeElapsedColumn(),
            TimeRemainingColumn(),
        )
        
        overall_task = progress.add_task("[bold magenta]Overall Progress", total=len(files))
        current_task = progress.add_task("[cyan]Starting...", total=None)

        layout = Layout()
        layout.split_column(
            Layout(name="header", size=3),
            Layout(name="main"),
            Layout(name="footer", size=10)
        )
        layout["header"].update(Panel("[bold white]RAG Ingestion Dashboard", style="on blue", box=box.ROUNDED))
        layout["main"].update(Panel(progress, title="Active Tasks", border_style="dim"))
        
        with Live(layout, refresh_per_second=10):
            for file_path in files:
                status = self.process_file(file_path, progress, current_task)
                progress.advance(overall_task)
                layout["footer"].update(Panel(self.make_stats_table(), title="Live Statistics", border_style="cyan"))
        
        self.print_summary()

    def print_summary(self):
        elapsed = _now() - self.stats["start_time"]
        summary = Table.grid(expand=True)
        summary.add_column(style="bold")
        summary.add_column(justify="right")
        summary.add_row("Success Rate", f"{(self.stats['processed']/(self.stats['processed']+self.stats['failed']+0.001))*100:.1f}%")
        summary.add_row("Avg Chunk Size", f"{self.stats['chunks_added']/(self.stats['processed']+0.001):.1f} per doc")
        summary.add_row("Total Processing Time", f"{elapsed:.1f}s")

        console.print("\n", Panel(summary, title="Ingestion Summary", border_style="green", expand=False))

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--source", default=str(DEFAULT_SOURCE_DIR))
    parser.add_argument("--chroma-dir", default=str(DEFAULT_CHROMA_DIR))
    parser.add_argument("--collection", default=DEFAULT_COLLECTION)
    parser.add_argument("--force", action="store_true")
    parser.add_argument("--no-tables", action="store_true")
    parser.add_argument("--no-image-markers", action="store_true")
    parser.add_argument("--chunk-chars", type=int, default=DEFAULT_CHUNK_CHARS)
    parser.add_argument("--overlap-chars", type=int, default=DEFAULT_CHUNK_OVERLAP_CHARS)
    parser.add_argument("--min-chunk-chars", type=int, default=DEFAULT_MIN_CHUNK_CHARS)
    parser.add_argument("--embed-model", default=DEFAULT_EMBED_MODEL)
    parser.add_argument("--embed-batch", type=int, default=DEFAULT_EMBED_BATCH)
    parser.add_argument("--test", action="store_true")
    args = parser.parse_args()

    processor = RagDocumentProcessor(
        source_dir=Path(args.source), chroma_dir=Path(args.chroma_dir),
        collection=args.collection, force_reprocess=args.force,
        include_tables=not args.no_tables, include_images_markers=not args.no_image_markers,
        chunk_chars=args.chunk_chars, overlap_chars=args.overlap_chars,
        min_chunk_chars=args.min_chunk_chars, embed_model=args.embed_model,
        embed_batch=args.embed_batch
    )

    processor.initialize()
    if args.test:
        files = []
        for ext in ["*.pdf", "*.docx", "*.pptx"]:
            files.extend([p for p in Path(args.source).rglob(ext) if not p.name.startswith(".")])
        if files:
            p = Progress(); t = p.add_task("Test"); processor.process_file(files[0], p, t)
            processor.print_summary()
        sys.exit(0)

    processor.process_all()

if __name__ == "__main__":
    try:
        main()
    except KeyboardInterrupt:
        console.print("\n[bold red]Interrupted by user.")
        sys.exit(130)